"""Generate a PowerPoint presentation of the Agrocosmos portal aimed at
government / regional administration audiences.

Run::

    py scripts/build_presentation.py

Output: ``presentation/agrocosmos_gov.pptx``.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from PIL import Image  # type: ignore  # bundled by python-pptx via Pillow
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt, Emu


# ── Brand palette ────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0E, 0x2B, 0x2A)        # deep teal — title backgrounds
ACCENT = RGBColor(0x2E, 0xA3, 0x7E)         # green agro accent
ACCENT_DARK = RGBColor(0x1E, 0x6E, 0x55)
LIGHT_BG = RGBColor(0xF5, 0xF8, 0xF6)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x14, 0x21, 0x20)
MUTED = RGBColor(0x55, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN = RGBColor(0xCC, 0x55, 0x33)


SLIDE_W_CM = 33.867   # 16:9 widescreen default in python-pptx (≈13.33 in)
SLIDE_H_CM = 19.05

SCREENS = Path(__file__).resolve().parent.parent / 'presentation' / 'screens'
ASSETS = Path(__file__).resolve().parent.parent / 'presentation' / 'assets'
LOGO = ASSETS / 'geo_app_logo.png'


# ── Helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = ACCENT
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill_rgb):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT,
                bullet_color=ACCENT, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        # bullet glyph as colored run
        bullet = p.add_run()
        bullet.text = '■  '
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        bullet.font.name = font
        # body
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = font
    return tb


def add_header_band(slide, title, subtitle=None):
    """Top band with title; used on content slides."""
    band_h = Cm(2.6)
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), band_h, DARK_BG)
    # Accent stripe
    add_rect(slide, 0, band_h, Cm(SLIDE_W_CM), Cm(0.18), ACCENT)
    add_text(
        slide, Cm(1.2), Cm(0.55), Cm(SLIDE_W_CM - 7.5), Cm(1.5),
        title, size=26, bold=True, color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(
            slide, Cm(1.2), Cm(1.55), Cm(SLIDE_W_CM - 7.5), Cm(0.9),
            subtitle, size=12, color=RGBColor(0xC2, 0xD8, 0xCF),
        )
    # Top-right GEO APP logo + brand strip
    if LOGO.exists():
        # Logo is square-ish; fit ~1.9 cm height, centered in band
        logo_h = Cm(1.9)
        logo_w = Cm(1.9)
        slide.shapes.add_picture(
            str(LOGO),
            Cm(SLIDE_W_CM - 2.5), Cm(0.35),
            width=logo_w, height=logo_h,
        )
        # Mini brand caption to the left of the logo
        add_text(
            slide, Cm(SLIDE_W_CM - 8.5), Cm(0.55), Cm(5.5), Cm(1.5),
            'Агрокосмос · ГЕО АП',
            size=11, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, Cm(SLIDE_W_CM - 8.5), Cm(1.45), Cm(5.5), Cm(1),
            'edunabazar.ru/agrocosmos',
            size=9, color=RGBColor(0xB6, 0xC8, 0xC2),
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )
    else:
        add_text(
            slide, Cm(SLIDE_W_CM - 6.5), Cm(0.6), Cm(5.5), Cm(1.4),
            'Агрокосмос  ·  edunabazar.ru',
            size=11, color=ACCENT, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def add_footer(slide, page_no, total):
    add_rect(slide, 0, Cm(SLIDE_H_CM - 0.7), Cm(SLIDE_W_CM), Cm(0.7), DARK_BG)
    add_text(
        slide, Cm(1.0), Cm(SLIDE_H_CM - 0.7), Cm(SLIDE_W_CM - 2), Cm(0.7),
        f'edunabazar.ru/agrocosmos     ·     слайд {page_no} / {total}',
        size=10, color=RGBColor(0xB6, 0xC8, 0xC2),
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )


def add_card(slide, x, y, w, h, title, body, *,
             title_color=ACCENT_DARK, accent=ACCENT):
    add_round(slide, x, y, w, h, CARD_BG)
    # left accent bar
    add_rect(slide, x, y, Cm(0.18), h, accent)
    add_text(
        slide, x + Cm(0.5), y + Cm(0.25), w - Cm(0.7), Cm(1.0),
        title, size=14, bold=True, color=title_color,
    )
    add_text(
        slide, x + Cm(0.5), y + Cm(1.05), w - Cm(0.7), h - Cm(1.2),
        body, size=11, color=TEXT,
    )


def add_screenshot(slide, img_path, x, y, w, h, *, caption=None):
    """Insert an image fitted into the (w x h) box, preserving aspect ratio,
    centered, with a thin frame and an optional caption underneath."""
    img_path = str(img_path)
    with Image.open(img_path) as im:
        iw, ih = im.size
    # Fit into box
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        # image wider — fit width
        nw = w
        nh = int(w / img_ar)
    else:
        nh = h
        nw = int(h * img_ar)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    # Frame background (subtle)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, nx - Emu(20000), ny - Emu(20000),
        nw + Emu(40000), nh + Emu(40000),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(0xCC, 0xD8, 0xD3)
    frame.line.width = Pt(0.5)
    frame.shadow.inherit = False
    slide.shapes.add_picture(img_path, nx, ny, width=nw, height=nh)
    if caption:
        add_text(
            slide, x, y + h + Cm(0.1), w, Cm(0.7),
            caption, size=10, color=MUTED, align=PP_ALIGN.CENTER,
        )


def add_metric(slide, x, y, w, h, value, label):
    add_round(slide, x, y, w, h, CARD_BG)
    add_text(
        slide, x, y + Cm(0.4), w, Cm(1.8),
        value, size=34, bold=True, color=ACCENT_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, x, y + h - Cm(1.1), w, Cm(0.9),
        label, size=11, color=MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# ── Slides ───────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Full-bleed dark background
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), DARK_BG)
    # Decorative accent bar
    add_rect(slide, 0, Cm(SLIDE_H_CM / 2 - 0.05),
             Cm(SLIDE_W_CM), Cm(0.1), ACCENT)
    # Eyebrow
    add_text(
        slide, Cm(2), Cm(3), Cm(SLIDE_W_CM - 4), Cm(1),
        'ЦИФРОВОЙ СПУТНИКОВЫЙ МОНИТОРИНГ ПОСЕВОВ',
        size=14, bold=True, color=ACCENT,
    )
    # Big title
    add_text(
        slide, Cm(2), Cm(4.2), Cm(SLIDE_W_CM - 4), Cm(4),
        'Агрокосмос',
        size=88, bold=True, color=WHITE,
    )
    add_text(
        slide, Cm(2), Cm(8), Cm(SLIDE_W_CM - 4), Cm(2.5),
        'Бесплатная геоинформационная система мониторинга\n'
        'вегетации сельхозугодий по всей территории России',
        size=22, color=RGBColor(0xC2, 0xD8, 0xCF),
    )
    # Subtitle / audience cue
    add_text(
        slide, Cm(2), Cm(13.5), Cm(SLIDE_W_CM - 4), Cm(1),
        'Презентация для региональных министерств сельского хозяйства\n'
        'и муниципальных администраций',
        size=14, color=RGBColor(0x9D, 0xB6, 0xAE),
    )
    # Footer brand + company logo
    add_text(
        slide, Cm(2), Cm(SLIDE_H_CM - 2.4), Cm(SLIDE_W_CM - 4), Cm(0.8),
        'edunabazar.ru/agrocosmos',
        size=14, bold=True, color=ACCENT,
    )
    add_text(
        slide, Cm(2), Cm(SLIDE_H_CM - 1.5), Cm(SLIDE_W_CM - 4), Cm(0.8),
        'Разработчик · ООО «ГЕО АП» · георешения для бизнеса',
        size=11, color=RGBColor(0x9D, 0xB6, 0xAE),
    )
    if LOGO.exists():
        # Logo top-right of title slide, medium size
        slide.shapes.add_picture(
            str(LOGO),
            Cm(SLIDE_W_CM - 5.5), Cm(1.5),
            width=Cm(3.8), height=Cm(3.8),
        )
    return slide


def slide_problem(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Проблема',
        'Региональные органы управления АПК сегодня оперируют отчётностью в Excel '
        'с лагом 2–4 недели и не имеют инструмента для оперативного контроля.',
    )

    items = [
        (
            'Отчётность приходит постфактум',
            'Сводки от хозяйств — раз в декаду или месяц, бумажные. К моменту '
            'выявления засухи или вымерзания пострадавшие площади уже потеряны.',
        ),
        (
            'Объективные данные — только у крупных холдингов',
            'Коммерческие сервисы (Sentinel Hub, Planet) стоят от 300 ₽/га/год. '
            'Малые хозяйства и муниципалитеты их позволить не могут.',
        ),
        (
            'Нет единой региональной картины',
            'Данные разбросаны по десяткам систем: Минсельхоз, Росреестр, '
            'Росгидромет, страховщики. Свести их в одну карту вручную — невозможно.',
        ),
        (
            'Аномалии замечают по жалобам, а не превентивно',
            'Засуха, паводок, поражение вредителями — становятся известны '
            'администрации, когда фермер уже подал заявление об ущербе.',
        ),
    ]
    x0 = Cm(1.5)
    y0 = Cm(3.6)
    card_w = Cm(SLIDE_W_CM - 3)
    card_h = Cm(2.8)
    gap = Cm(0.4)
    for i, (t, b) in enumerate(items):
        add_card(slide, x0, y0 + i * (card_h + gap),
                 card_w, card_h, t, b,
                 title_color=WARN, accent=WARN)
    add_footer(slide, page, total)


def slide_solution(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Решение',
        'Агрокосмос — единая ГИС-платформа спутникового мониторинга вегетации, '
        'бесплатная для конечного пользователя и публично размещённая.',
    )

    # Left: pitch text
    add_text(
        slide, Cm(1.5), Cm(3.6), Cm(15), Cm(1.2),
        'Что мы даём региону:',
        size=18, bold=True, color=ACCENT_DARK,
    )
    bullets = [
        'Карта-светофор всех муниципальных районов России — '
        'оперативное состояние посевов в одном экране.',
        'Архив NDVI с 2000 года — динамика вегетации за 25+ лет '
        'с наложением многолетней нормы.',
        'Автоматические оповещения об аномалиях — отклонение от '
        'базовой линии и резкое падение индекса.',
        'Drill-down: страна → регион → район → конкретное поле '
        '(до 10-метрового разрешения).',
        'Интеграция через открытое API — данные доступны другим '
        'информационным системам региона.',
    ]
    add_bullets(slide, Cm(1.5), Cm(4.6), Cm(15), Cm(10), bullets, size=14)

    # Right: key metrics
    rx = Cm(17.5)
    rw = Cm(SLIDE_W_CM - 19)
    rh = Cm(3.3)
    metrics = [
        ('2 300+', 'муниципальных районов покрыто'),
        ('25 лет', 'архив MODIS NDVI с 2000 года'),
        ('10 м', 'разрешение Sentinel-2 для полей'),
        ('0 ₽', 'для конечного пользователя'),
    ]
    for i, (v, lab) in enumerate(metrics):
        add_metric(slide, rx, Cm(3.6) + i * (rh + Cm(0.3)),
                   rw, rh, v, lab)

    add_footer(slide, page, total)


def slide_coverage(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Покрытие и источники данных',
        'Открытые спутниковые программы NASA/USGS и ESA — никакой зависимости '
        'от коммерческих провайдеров.',
    )
    sources = [
        (
            'MODIS Terra+Aqua (250 м)',
            'Архив с 2000 года. 16-дневные NDVI-композиты обновляются '
            'ежедневно. Используется как многолетняя норма (baseline) '
            'и для всероссийской карты-светофора.',
        ),
        (
            'Sentinel-2 L2A (10 м)',
            'Программа ESA. Композиты раз в ~5 дней. Используется для '
            'полевого мониторинга — детализация до контура поля.',
        ),
        (
            'Landsat 8/9 C2L2 (30 м)',
            'Программа NASA/USGS. 16-дневный цикл. Подстраховка для '
            'периодов плотной облачности и для калибровки HLS-fused.',
        ),
        (
            'HLS Fusion (гармонизация)',
            'Собственный пайплайн склеивает Sentinel-2 и Landsat в единый '
            'плотный временной ряд — закрывает пропуски от облачности.',
        ),
    ]
    x0 = Cm(1.5)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 3 - 0.6) / 2)
    card_h = Cm(5.5)
    for i, (t, b) in enumerate(sources):
        col = i % 2
        row = i // 2
        x = x0 + col * (card_w + Cm(0.6))
        y = y0 + row * (card_h + Cm(0.5))
        add_card(slide, x, y, card_w, card_h, t, b)

    add_footer(slide, page, total)


def slide_what_official_sees(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Что видит сотрудник министерства',
        'Иерархия экранов: страна → регион → район → поле. '
        'От обзора до детали — три клика.',
    )

    steps = [
        ('1', 'Карта России',
         'Хороплет 2 300+ районов раскрашен по NDVI z-score. '
         'Красный = режим дефицита, зелёный = выше нормы. '
         'Аномалии регионального масштаба видны за один взгляд.'),
        ('2', 'Регион → район',
         'Клик по области открывает её муниципалитеты с тем же '
         'светофором. Можно сравнить районы между собой и сразу '
         'выявить отстающие.'),
        ('3', 'NDVI-график района',
         'Текущий сезон на фоне многолетнего коридора нормы. '
         'Видно, идёт ли вегетация по графику или отстаёт; '
         'сразу подсвечены даты, когда отклонение превысило 1.5σ.'),
        ('4', 'Конкретное поле',
         'По полигонам Россельхозбанка / ЕСФД доступна история NDVI '
         'на разрешении 10 м, фенологические метрики (SOS/POS/EOS) '
         'и экспорт в CSV для отчётности.'),
    ]
    x0 = Cm(1.2)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 2.4 - 0.9) / 4)
    card_h = Cm(11.3)
    for i, (n, t, b) in enumerate(steps):
        x = x0 + i * (card_w + Cm(0.3))
        add_round(slide, x, y0, card_w, card_h, CARD_BG)
        # Step number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Cm(0.5), y0 + Cm(0.5),
            Cm(1.4), Cm(1.4),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        add_text(
            slide, x + Cm(0.5), y0 + Cm(0.5), Cm(1.4), Cm(1.4),
            n, size=22, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, x + Cm(0.4), y0 + Cm(2.2), card_w - Cm(0.8), Cm(1.3),
            t, size=15, bold=True, color=ACCENT_DARK,
        )
        add_text(
            slide, x + Cm(0.4), y0 + Cm(3.6), card_w - Cm(0.8),
            card_h - Cm(4.0), b, size=11, color=TEXT,
        )

    add_footer(slide, page, total)


def slide_screen_country(prs, page, total):
    """Big screenshot: all-Russia choropleth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Карта-светофор всей страны',
        'Хороплет 2 300+ муниципальных районов в одном экране. Цвет — '
        'отношение текущего NDVI к многолетней норме (z-score).',
    )

    # Screenshot fills most of the slide
    img = SCREENS / '01_dashboard_all_russia.png'
    add_screenshot(
        slide, img,
        Cm(1.5), Cm(3.4), Cm(SLIDE_W_CM - 12), Cm(14.5),
    )

    # Caption legend on right
    rx = Cm(SLIDE_W_CM - 10)
    rw = Cm(8.5)
    add_text(
        slide, rx, Cm(3.6), rw, Cm(1),
        'Что показано',
        size=18, bold=True, color=ACCENT_DARK,
    )
    items = [
        'Каждый район окрашен в один из 7 классов NDVI.',
        'Красный/оранжевый = режим дефицита (засуха, отставание).',
        'Зелёный = в норме или выше многолетнего среднего.',
        'Серый = недостаточно данных за выбранный период.',
        'Слайдер «Состояние на …» — выбор даты композита.',
        'Один клик по району — переход в детальный экран.',
    ]
    add_bullets(slide, rx, Cm(4.7), rw, Cm(11), items, size=12)

    add_text(
        slide, rx, Cm(15.5), rw, Cm(2),
        'Производительность: первая отрисовка ~3 секунды на холодный кеш, '
        '<200 мс на тёплый. Размер payload снижен с 25 МБ до 4 МБ за счёт '
        'серверной геометрической генерализации.',
        size=10, color=MUTED,
    )

    add_footer(slide, page, total)


def slide_screen_region(prs, page, total):
    """Region drill-down with NDVI dynamics chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Регион и динамика NDVI',
        'Клик по субъекту открывает его районы с тем же светофором '
        'и график NDVI за текущий сезон на фоне многолетней нормы.',
    )

    img = SCREENS / '02_region_ndvi_chart.png'
    add_screenshot(
        slide, img,
        Cm(1.5), Cm(3.4), Cm(SLIDE_W_CM - 3), Cm(14.5),
        caption='Республика Крым: муниципалитеты раскрашены по NDVI; '
                'в правом верхнем углу — текущий сезон vs многолетняя норма',
    )

    add_footer(slide, page, total)


def slide_screen_report(prs, page, total):
    """Region NDVI report page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Готовые отчёты по региону и району',
        'Отдельная страница «Отчёт NDVI» — для печати, рассылки в кризис-штаб '
        'или прикрепления к министерским сводкам.',
    )

    # Left: screenshot
    add_screenshot(
        slide, SCREENS / '03_report_region.png',
        Cm(1.5), Cm(3.4), Cm(SLIDE_W_CM - 12), Cm(14.5),
    )

    # Right: explanation
    rx = Cm(SLIDE_W_CM - 10)
    rw = Cm(8.5)
    add_text(
        slide, rx, Cm(3.6), rw, Cm(1),
        'Что в отчёте',
        size=18, bold=True, color=ACCENT_DARK,
    )
    items = [
        'Сводные показатели региона: текущее значение NDVI, число районов '
        'с данными, среднее значение по медиане.',
        'Многолетний тренд: текущий сезон vs медиана за 20+ лет архива.',
        'Разбивка по районам с маркировкой аномалий (z-score > 1.5σ).',
        'Кнопки «Печать» / «E-mail» / «Telegram» — оперативная рассылка.',
        'Тот же отчёт можно построить на уровне отдельного района.',
    ]
    add_bullets(slide, rx, Cm(4.7), rw, Cm(11), items, size=12)

    add_text(
        slide, rx, Cm(15.5), rw, Cm(2),
        'Примеры запросов: «среднее NDVI по Республике Крым на дату X», '
        '«районы Республики Крым с отклонением > 1.5σ», '
        '«фенология Бахчисарайского района за 2026 год».',
        size=10, color=MUTED,
    )

    add_footer(slide, page, total)


def slide_screen_raster(prs, page, total):
    """Detailed raster dashboard with NDVI overlay."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Растровый мониторинг угодий',
        'Sentinel-2 (10 м) и Landsat 8/9 (30 м) — псевдоцветный NDVI поверх '
        'границ района. Слой включается одной галкой в боковой панели.',
    )

    add_screenshot(
        slide, SCREENS / '04_raster_dashboard.png',
        Cm(1.5), Cm(3.4), Cm(SLIDE_W_CM - 3), Cm(14.5),
        caption='Растровый дашборд: выбор сенсора (S2/L8), даты композита, '
                'прозрачности; легенда NDVI 0.0–1.0 в боковой панели',
    )

    add_footer(slide, page, total)


def slide_alerts(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Раннее обнаружение проблем',
        'Система ежедневно сравнивает фактический NDVI с многолетней нормой '
        'и сама присылает уведомление, если что-то пошло не так.',
    )

    # Two-column layout: types of alerts + workflow
    add_text(
        slide, Cm(1.5), Cm(3.6), Cm(15), Cm(1),
        'Два типа автоматических тревог',
        size=18, bold=True, color=ACCENT_DARK,
    )
    add_card(
        slide, Cm(1.5), Cm(4.7), Cm(15), Cm(3.5),
        'Отклонение от базовой линии (baseline drift)',
        'Устойчивое снижение NDVI более чем на 1.5 стандартных отклонения '
        'от среднего значения за тот же день года за последние 20+ лет. '
        'Признак затяжной засухи, переувлажнения или поражения посева.',
    )
    add_card(
        slide, Cm(1.5), Cm(8.4), Cm(15), Cm(3.5),
        'Резкое падение (sudden drop)',
        'Скачкообразное снижение между двумя соседними 16-дневными '
        'композитами. Сигнализирует о точечных событиях: град, '
        'возвратные заморозки, локальный пожар, нашествие вредителей.',
    )

    # Right: how delivered
    add_text(
        slide, Cm(17.5), Cm(3.6), Cm(SLIDE_W_CM - 19), Cm(1),
        'Как доставляются',
        size=18, bold=True, color=ACCENT_DARK,
    )
    bullets = [
        'E-mail подписчику с прямой ссылкой на район/поле.',
        'Подписка настраивается по списку регионов и/или районов.',
        'Отдельные галки: «аномалии» и «регулярные сводки».',
        'Доступ через личный кабинет — без регистрации.',
        'Планируется: SMS, Telegram-бот, web-push.',
    ]
    add_bullets(slide, Cm(17.5), Cm(4.7), Cm(SLIDE_W_CM - 19), Cm(7.2),
                bullets, size=13)

    # Bottom: example threshold annotation
    add_round(slide, Cm(1.5), Cm(12.3), Cm(SLIDE_W_CM - 3), Cm(3),
              RGBColor(0xE9, 0xF3, 0xEE))
    add_text(
        slide, Cm(2), Cm(12.5), Cm(SLIDE_W_CM - 4), Cm(0.8),
        'Пример практического применения',
        size=13, bold=True, color=ACCENT_DARK,
    )
    add_text(
        slide, Cm(2), Cm(13.2), Cm(SLIDE_W_CM - 4), Cm(2),
        'Региональный кризис-штаб получает утреннюю рассылку: '
        '«Краснодарский край — 4 района в зоне baseline drift на 28-й '
        'день после посева». До звонков от хозяйств остаётся 7–10 дней '
        '— время на полевую проверку, корректировку планов уборки и '
        'предупреждение страховщиков.',
        size=12, color=TEXT,
    )

    add_footer(slide, page, total)


def slide_use_cases(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Сценарии применения для региона',
        'Шесть готовых рабочих процессов, которые система закрывает «из коробки».',
    )
    cases = [
        ('Контроль посевной кампании',
         'Сравнение SOS (start of season) текущего года с многолетней нормой '
         'и соседними областями. Видны районы, где сев идёт с задержкой.'),
        ('Мониторинг засухи и переувлажнения',
         'Оперативный учёт районов в режиме дефицита или избытка влаги '
         'через NDVI z-score. Подкреплено архивом аномалий за 25 лет.'),
        ('Подготовка к чрезвычайным ситуациям',
         'Картирование посевов в зоне риска после града, заморозков, '
         'пожаров. Объективная оценка пострадавших площадей в гектарах.'),
        ('Поддержка страховых выплат',
         'Независимый источник данных для подтверждения наступления '
         'страхового случая (засуха, недобор урожая) — сокращает споры.'),
        ('Учёт неиспользуемых сельхозземель',
         'Поля с устойчиво низким NDVI многолетне = заброшены или выведены. '
         'Помощь Россельхознадзору и Росреестру в инвентаризации.'),
        ('Открытый доступ для фермеров',
         'Региональные хозяйства получают тот же инструмент бесплатно. '
         'Снижается нагрузка на «горячие линии» министерства.'),
    ]
    x0 = Cm(1.0)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 2.0 - 1.2) / 3)
    card_h = Cm(5.5)
    for i, (t, b) in enumerate(cases):
        col = i % 3
        row = i // 3
        x = x0 + col * (card_w + Cm(0.6))
        y = y0 + row * (card_h + Cm(0.5))
        add_card(slide, x, y, card_w, card_h, t, b)

    add_footer(slide, page, total)


def slide_methodology(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Методология и научная база',
        'Стандартные подходы дистанционного зондирования, прозрачные формулы, '
        'воспроизводимые результаты.',
    )
    items = [
        ('NDVI — индекс вегетации',
         'Normalized Difference Vegetation Index = (NIR − RED) / (NIR + RED). '
         'Универсальная мера фотосинтетической активности. Используется '
         'Минсельхозом РФ, ФАО, ESA, NASA как опорный показатель.'),
        ('Зональная статистика по полигонам',
         'Для каждого поля считается area-weighted среднее NDVI с учётом '
         'облачности и качества пикселя. Облачные пиксели отсекаются '
         'через QA-маски Sentinel-2 SCL и Landsat QA_PIXEL.'),
        ('Многолетняя норма (baseline)',
         'По каждому дню года агрегируется среднее NDVI и std за 20+ лет '
         'архива MODIS. Текущее значение сравнивается через z-score.'),
        ('Фенологические метрики',
         'Алгоритм TIMESAT-like: SOS (start of season), POS (peak of season), '
         'EOS (end of season), LOS (length of season). Помогает оценить '
         'смещение фенофаз относительно нормы.'),
        ('Сглаживание и детекция выбросов',
         'Savitzky-Golay фильтр + правило 3σ на остатках убирают '
         'облачные/снежные выбросы из временного ряда без потери '
         'реальных аномалий.'),
        ('Гармонизация S2 и Landsat (HLS)',
         'Единый пайплайн склеивает разнородные снимки в один временной ряд '
         '— увеличивает плотность наблюдений в облачные периоды в 2–3 раза.'),
    ]
    x0 = Cm(1.0)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 2.0 - 1.2) / 3)
    card_h = Cm(5.5)
    for i, (t, b) in enumerate(items):
        col = i % 3
        row = i // 3
        x = x0 + col * (card_w + Cm(0.6))
        y = y0 + row * (card_h + Cm(0.5))
        add_card(slide, x, y, card_w, card_h, t, b)

    add_footer(slide, page, total)


def slide_tech(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Технологическая основа',
        'Импортонезависимый стек на открытых компонентах. Российские серверы.',
    )

    # Stack columns
    cols = [
        ('Бэкенд',
         ['Python 3.12, Django 5',
          'GeoDjango для геоопераций',
          'PostgreSQL 16 + PostGIS 3.4',
          'Redis для кеша и сессий',
          'Gunicorn + Nginx']),
        ('Данные',
         ['MODIS / Sentinel-2 / Landsat',
          'GeoTIFF, COG-форматы',
          'Зональная статистика на лету',
          'Pre-aggregate за 25 лет',
          'Mapbox Vector Tiles']),
        ('Инфраструктура',
         ['Docker Compose (прод)',
          'GitHub Actions CI/CD',
          'Ежедневные бэкапы PostgreSQL',
          'PVE-шлюз + 2 VM',
          'Российский хостинг']),
        ('Доступы и безопасность',
         ['HTTPS (Let\'s Encrypt)',
          'Rate limiting на API',
          'Личный кабинет с подписками',
          'Открытое REST-API',
          'Аудит-логи действий']),
    ]
    x0 = Cm(1.2)
    y0 = Cm(3.6)
    col_w = Cm((SLIDE_W_CM - 2.4 - 1.2) / 4)
    col_h = Cm(11.3)
    for i, (t, items) in enumerate(cols):
        x = x0 + i * (col_w + Cm(0.4))
        add_round(slide, x, y0, col_w, col_h, CARD_BG)
        add_rect(slide, x, y0, col_w, Cm(1.3), ACCENT_DARK)
        add_text(
            slide, x, y0, col_w, Cm(1.3),
            t, size=15, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # bullets
        add_bullets(
            slide, x + Cm(0.4), y0 + Cm(1.6),
            col_w - Cm(0.8), col_h - Cm(1.8),
            items, size=11,
        )

    add_footer(slide, page, total)


def slide_integration(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Интеграция с региональными системами',
        'Открытое REST-API позволяет встроить данные Агрокосмоса в существующие '
        'информационные контуры министерств и муниципалитетов.',
    )
    cards = [
        ('Открытое REST-API',
         'JSON-эндпоинты для получения NDVI по региону, району или полю '
         'с любым временным окном. Документация — `docs/AGROCOSMOS_API.md`. '
         'Rate-limit: до 300 запросов в минуту с IP.'),
        ('GeoJSON / Vector Tiles',
         'Полигоны районов и полей отдаются в стандартных форматах для '
         'импорта в QGIS, ArcGIS и любые региональные ГИС. MVT-тайлы '
         'для интеграции с веб-картой региона.'),
        ('E-mail подписки',
         'Подписки можно настраивать программно через API: список '
         'регионов/районов, типы оповещений. Подходит для подключения '
         'кризис-штаба или служб госмониторинга.'),
        ('Импорт пользовательских полигонов',
         'Загрузка собственных контуров (Россельхозбанк, ЕСФД, '
         'кадастровые карты) — система считает NDVI по любому полигону '
         'на лету, без предварительной разметки.'),
    ]
    x0 = Cm(1.5)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 3 - 0.6) / 2)
    card_h = Cm(5.5)
    for i, (t, b) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = x0 + col * (card_w + Cm(0.6))
        y = y0 + row * (card_h + Cm(0.5))
        add_card(slide, x, y, card_w, card_h, t, b)

    add_footer(slide, page, total)


def slide_economics(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Экономика для региона',
        'Сколько стоила бы аналогичная задача через коммерческого провайдера '
        'и почему мы делаем это бесплатно.',
    )

    # Comparison table
    rows = [
        ('Параметр', 'Коммерческие сервисы', 'Агрокосмос'),
        ('Стоимость для конечного пользователя',
         'от 300 ₽ за гектар в год', 'бесплатно'),
        ('Минимальный пакет',
         'обычно от 1 000 га', 'без минимума'),
        ('Архив NDVI', 'до 5–10 лет', '25 лет (с 2000)'),
        ('Покрытие России',
         'выборочно, по подписке', 'все 2 300+ районов'),
        ('Открытое API',
         'на дополнительной лицензии', 'входит'),
        ('Зависимость от иностранных платформ',
         'высокая (Sentinel Hub, Planet)',
         'данные открытые, обработка своя'),
    ]
    table_x = Cm(1.5)
    table_y = Cm(3.6)
    table_w = Cm(SLIDE_W_CM - 3)
    row_h = Cm(1.2)
    col_widths = [Cm((SLIDE_W_CM - 3) * 0.34),
                  Cm((SLIDE_W_CM - 3) * 0.33),
                  Cm((SLIDE_W_CM - 3) * 0.33)]
    for ri, row in enumerate(rows):
        is_header = ri == 0
        # row background
        bg = ACCENT_DARK if is_header else (CARD_BG if ri % 2 else RGBColor(0xEC, 0xF4, 0xEF))
        add_rect(slide, table_x, table_y + ri * row_h,
                 table_w, row_h, bg)
        cx = table_x
        for ci, val in enumerate(row):
            color = WHITE if is_header else TEXT
            bold = is_header or ci == 2
            add_text(
                slide, cx + Cm(0.3), table_y + ri * row_h,
                col_widths[ci] - Cm(0.6), row_h,
                val, size=12, bold=bold, color=color,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += col_widths[ci]

    # Bottom note
    note_y = table_y + len(rows) * row_h + Cm(0.6)
    add_round(slide, Cm(1.5), note_y, Cm(SLIDE_W_CM - 3), Cm(2.4),
              RGBColor(0xE9, 0xF3, 0xEE))
    add_text(
        slide, Cm(2), note_y + Cm(0.3), Cm(SLIDE_W_CM - 4), Cm(0.8),
        'Почему бесплатно',
        size=14, bold=True, color=ACCENT_DARK,
    )
    add_text(
        slide, Cm(2), note_y + Cm(1.0), Cm(SLIDE_W_CM - 4), Cm(1.4),
        'Спутниковые данные NASA/USGS и ESA публичны и бесплатны для всех. '
        'Мы платим только за серверы и каналы. Сервис финансируется через '
        'смежные платные продукты «Единого Базара» (маркетплейс агропродукции), '
        'и поэтому может оставаться бесплатным для региональных органов и фермеров.',
        size=12, color=TEXT,
    )

    add_footer(slide, page, total)


def slide_pilot(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Пилот для региона',
        'Что мы можем сделать вместе за один сельхоз-сезон.',
    )
    phases = [
        ('Месяц 1 — Подключение',
         'Импорт ваших полигонов полей (Россельхозбанк, ЕСФД, кадастр). '
         'Выдача доступов специалистам министерства. Настройка '
         'еженедельной email-сводки.'),
        ('Месяц 2–3 — Калибровка',
         'Сравнение прогнозов NDVI с фактическими отчётами хозяйств. '
         'Настройка пороговых значений тревог под особенности региона. '
         'Обучение операторов кризис-штаба.'),
        ('Месяц 4–6 — Эксплуатация',
         'Полноценный мониторинг посевной и вегетационного сезона. '
         'Еженедельные сводки по районам. Отработка автоматических '
         'оповещений об аномалиях.'),
        ('Месяц 7–9 — Урожай',
         'Сравнение прогноза по фенологии с фактическим сроком уборки. '
         'Картирование пострадавших участков. Помощь в подготовке '
         'отчётности по итогам года.'),
        ('Месяц 10–12 — Оценка',
         'Совместный анализ: сколько раз система предсказала проблему '
         'раньше отчётности, экономический эффект, доработка под '
         'специфику региона.'),
    ]
    y0 = Cm(3.6)
    row_h = Cm(2.4)
    for i, (t, b) in enumerate(phases):
        y = y0 + i * (row_h + Cm(0.2))
        # Number badge
        badge_x = Cm(1.5)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, badge_x, y + Cm(0.3),
            Cm(1.7), Cm(1.7),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        add_text(
            slide, badge_x, y + Cm(0.3), Cm(1.7), Cm(1.7),
            str(i + 1), size=22, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Card
        cx = badge_x + Cm(2.2)
        cw = Cm(SLIDE_W_CM) - cx - Cm(1.5)
        add_round(slide, cx, y, cw, row_h, CARD_BG)
        add_rect(slide, cx, y, Cm(0.18), row_h, ACCENT)
        add_text(
            slide, cx + Cm(0.5), y + Cm(0.25), cw - Cm(0.7), Cm(0.9),
            t, size=14, bold=True, color=ACCENT_DARK,
        )
        add_text(
            slide, cx + Cm(0.5), y + Cm(1.0), cw - Cm(0.7), row_h - Cm(1.2),
            b, size=11, color=TEXT,
        )

    add_footer(slide, page, total)


def slide_legal(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Соответствие требованиям',
        'Юридические и инфраструктурные аспекты для работы с госорганами.',
    )
    cards = [
        ('Хранение данных в РФ',
         'Все серверы — Proxmox VE на территории Российской Федерации. '
         'Резервное копирование PostgreSQL — на сетевое хранилище в РФ. '
         'Соответствие требованиям 152-ФЗ для персональных данных.'),
        ('Открытые источники',
         'Спутниковые данные NASA/USGS и ESA — публичные, без ограничений '
         'на использование, в том числе государственными органами РФ. '
         'Лицензии Public Domain / CC-BY.'),
        ('Импортонезависимый стек',
         'Python, PostgreSQL, Redis, Nginx — open-source. '
         'Все компоненты находятся в реестре отечественного ПО или имеют '
         'эквивалент. Нет зависимости от санкционных вендоров.'),
        ('Прозрачная методология',
         'Алгоритмы расчёта NDVI и аномалий полностью документированы. '
         'Тестовая сьюта на каждый коммит. Реализация воспроизводима '
         'и может быть проверена независимыми экспертами.'),
        ('Аудит и логирование',
         'Все действия пользователей и изменения подписок фиксируются. '
         'Доступ к административной панели — по белому списку логинов. '
         'HTTPS, rate limiting, защита от DDoS на уровне Nginx.'),
        ('Резервирование',
         'Ежедневные дампы PostgreSQL (30-дневное окно). Возможность '
         'переноса инсталляции в защищённый контур регионального '
         'министерства за 1–2 недели.'),
    ]
    x0 = Cm(1.0)
    y0 = Cm(3.6)
    card_w = Cm((SLIDE_W_CM - 2.0 - 1.2) / 3)
    card_h = Cm(5.5)
    for i, (t, b) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x0 + col * (card_w + Cm(0.6))
        y = y0 + row * (card_h + Cm(0.5))
        add_card(slide, x, y, card_w, card_h, t, b)

    add_footer(slide, page, total)


def slide_roadmap(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), LIGHT_BG)
    add_header_band(
        slide,
        'Дорожная карта',
        'Что уже работает и куда движемся в ближайшие 12 месяцев.',
    )

    # Two columns: now / next
    add_text(
        slide, Cm(1.5), Cm(3.6), Cm(SLIDE_W_CM / 2 - 2), Cm(1),
        'Уже в продакшене',
        size=20, bold=True, color=ACCENT_DARK,
    )
    now_items = [
        'Ежедневный пайплайн MODIS — все районы РФ',
        'Sentinel-2 + Landsat 8/9 для пользовательских полей',
        'HLS-fusion (склейка разнородных снимков)',
        'Карта-светофор по 2 300+ районам',
        'Личный кабинет с подписками и оповещениями',
        'Открытое REST-API + векторные тайлы',
        'PWA для мобильных устройств',
        '25-летний baseline и z-score',
    ]
    add_bullets(slide, Cm(1.5), Cm(4.7),
                Cm(SLIDE_W_CM / 2 - 2), Cm(11),
                now_items, size=13)

    # Roadmap column
    add_text(
        slide, Cm(SLIDE_W_CM / 2 + 0.5), Cm(3.6),
        Cm(SLIDE_W_CM / 2 - 2), Cm(1),
        'Планы (12 месяцев)',
        size=20, bold=True, color=ACCENT_DARK,
    )
    next_items = [
        'Telegram-бот с push-оповещениями',
        'Интеграция с ЕСФД и Россельхозбанком (полигоны)',
        'Прогноз урожайности по фенологии (ML-модель)',
        'Влагозапас в почве по SMAP/SAR',
        'Детекция возгораний по MODIS Active Fire',
        'Личные кабинеты для региональных министерств',
        'Экспорт отчётов в формате Минсельхоза РФ',
        'Подписка для страховщиков (API + dashboard)',
    ]
    add_bullets(slide, Cm(SLIDE_W_CM / 2 + 0.5), Cm(4.7),
                Cm(SLIDE_W_CM / 2 - 2), Cm(11),
                next_items, size=13, bullet_color=WARN)

    add_footer(slide, page, total)


def slide_call_to_action(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, Cm(SLIDE_W_CM), Cm(SLIDE_H_CM), DARK_BG)
    # Decorative accent
    add_rect(slide, 0, Cm(SLIDE_H_CM / 2 - 0.05),
             Cm(SLIDE_W_CM), Cm(0.1), ACCENT)

    add_text(
        slide, Cm(2), Cm(3), Cm(SLIDE_W_CM - 4), Cm(1),
        'СЛЕДУЮЩИЙ ШАГ',
        size=14, bold=True, color=ACCENT,
    )
    add_text(
        slide, Cm(2), Cm(4.2), Cm(SLIDE_W_CM - 4), Cm(3),
        'Запустим пилот\nдля вашего региона',
        size=58, bold=True, color=WHITE,
    )
    add_text(
        slide, Cm(2), Cm(9.5), Cm(SLIDE_W_CM - 4), Cm(2),
        'Подключение силами нашей команды, без затрат региона.\n'
        'Доступ министерству, обучение операторов, сопровождение.',
        size=18, color=RGBColor(0xC2, 0xD8, 0xCF),
    )

    # Contact card
    cx = Cm(2)
    cy = Cm(13.5)
    cw = Cm(SLIDE_W_CM - 4)
    ch = Cm(3.5)
    add_round(slide, cx, cy, cw, ch, RGBColor(0x1A, 0x3D, 0x3A))
    add_rect(slide, cx, cy, Cm(0.2), ch, ACCENT)
    add_text(
        slide, cx + Cm(0.7), cy + Cm(0.4), cw - Cm(1), Cm(0.9),
        'Контакты',
        size=14, bold=True, color=ACCENT,
    )
    add_text(
        slide, cx + Cm(0.7), cy + Cm(1.3), cw - Cm(1), Cm(0.9),
        'Сайт:        edunabazar.ru/agrocosmos',
        size=14, color=WHITE,
    )
    add_text(
        slide, cx + Cm(0.7), cy + Cm(2.0), cw - Cm(1), Cm(0.9),
        'Личный кабинет:  edunabazar.ru/me/agrocosmos',
        size=14, color=WHITE,
    )
    add_text(
        slide, cx + Cm(0.7), cy + Cm(2.7), cw - Cm(1), Cm(0.9),
        'E-mail:      info@edunabazar.ru',
        size=14, color=WHITE,
    )

    # Company line + logo (bottom-right)
    add_text(
        slide, Cm(2), Cm(SLIDE_H_CM - 1.0), Cm(SLIDE_W_CM - 4), Cm(0.7),
        'Разработка и сопровождение · ООО «ГЕО АП» · георешения для бизнеса',
        size=10, color=RGBColor(0x9D, 0xB6, 0xAE),
    )
    if LOGO.exists():
        slide.shapes.add_picture(
            str(LOGO),
            Cm(SLIDE_W_CM - 4.5), Cm(1.0),
            width=Cm(3.0), height=Cm(3.0),
        )


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    # 16:9
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_coverage,
        slide_what_official_sees,
        slide_screen_country,
        slide_screen_region,
        slide_screen_report,
        slide_screen_raster,
        slide_alerts,
        slide_use_cases,
        slide_methodology,
        slide_tech,
        slide_integration,
        slide_economics,
        slide_pilot,
        slide_legal,
        slide_roadmap,
        slide_call_to_action,
    ]
    total = len(builders)

    for idx, build in enumerate(builders, start=1):
        if build is slide_title:
            build(prs)
        else:
            build(prs, idx, total)

    out_dir = Path(__file__).resolve().parent.parent / 'presentation'
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / 'agrocosmos_gov.pptx'
    prs.save(out_path)
    print(f'Saved: {out_path}')


if __name__ == '__main__':
    main()
